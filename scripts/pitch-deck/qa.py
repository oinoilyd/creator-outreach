#!/usr/bin/env python3
"""
Programmatic QA for the rendered pitch deck.

Checks:
1. All 126 slides exist
2. Each slide has the expected brand mark + footer + topic tag
3. Text-fit heuristic — flag boxes where char count likely overflows shape
4. No empty/blank slides

Run: python3 scripts/pitch-deck/qa.py
"""
from pptx import Presentation
from pptx.util import Emu
import sys

PATH = "docs/pitch-deck.pptx"
EXPECTED_SLIDES = 126

# Heuristic: at 13pt Calibri body, ~11 chars per inch wide × ~4.5 lines per inch tall
# = ~50 chars per square inch. Scales as (13/pt)².
# Padded with 0.9 safety factor (we want to flag *before* actual overflow).
def estimate_fit(text, w_in, h_in, font_pt):
    if not text or not font_pt:
        return True, 0, 0
    base = 50.0  # chars per sq-inch at 13pt
    capacity = w_in * h_in * base * (13.0 / font_pt) ** 2 * 0.9
    return len(text) <= capacity, len(text), capacity


def main():
    prs = Presentation(PATH)
    n_slides = len(prs.slides)
    issues = []
    warnings = []

    print(f"Loaded: {PATH}")
    print(f"Slides: {n_slides} (expected {EXPECTED_SLIDES})")
    print(f"Slide size: {prs.slide_width / 914400:.2f} × {prs.slide_height / 914400:.2f} in")
    print()

    if n_slides != EXPECTED_SLIDES:
        issues.append(f"Slide count mismatch: {n_slides} ≠ {EXPECTED_SLIDES}")

    # Per-slide checks
    for i, slide in enumerate(prs.slides, start=1):
        slide_label = f"Slide {i:03d}"

        # Collect all text on the slide
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            all_text.append(run.text)
        full_text = " | ".join(all_text)

        # 1. Must contain "Creator Outreach" (brand wordmark)
        if "Creator Outreach" not in full_text and "creatoroutreach" not in full_text.lower():
            issues.append(f"{slide_label}: Missing 'Creator Outreach' wordmark")

        # 2. Must contain footer attribution
        if "Gaynor Media" not in full_text:
            issues.append(f"{slide_label}: Missing 'Gaynor Media' footer")

        # 3. Must contain slide number "NNN / 126"
        page_str = f"{i:03d} / 126"
        if page_str not in full_text:
            issues.append(f"{slide_label}: Missing page indicator '{page_str}'")

        # 4. Slide should have meaningful content (>200 chars of text)
        if len(full_text) < 200:
            warnings.append(f"{slide_label}: Only {len(full_text)} chars of text — may be sparse")

        # 5. Per-textbox fit check
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text or len(text) < 30:
                continue  # too small to overflow
            w_in = shape.width / 914400
            h_in = shape.height / 914400
            # Get the largest font size in the box
            max_pt = 0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_pt = max(max_pt, run.font.size.pt)
            if max_pt == 0:
                max_pt = 14
            fits, chars, capacity = estimate_fit(text, w_in, h_in, max_pt)
            if not fits:
                ratio = chars / capacity if capacity else 999
                # Only flag if substantially over (>1.3x), because the heuristic is loose
                if ratio > 1.3:
                    warnings.append(
                        f"{slide_label}: Text may overflow ({chars} chars in "
                        f"{w_in:.1f}×{h_in:.1f}in box @ {max_pt:.0f}pt, capacity ~{capacity:.0f}). "
                        f"Preview: {text[:60]!r}"
                    )

    # Report
    print(f"━━━ Issues ({len(issues)}) ━━━")
    for i in issues[:20]:
        print(f"  ✗ {i}")
    if len(issues) > 20:
        print(f"  … and {len(issues) - 20} more")

    print()
    print(f"━━━ Warnings ({len(warnings)}) ━━━")
    for w in warnings[:30]:
        print(f"  ! {w}")
    if len(warnings) > 30:
        print(f"  … and {len(warnings) - 30} more")

    print()
    if issues:
        print(f"✗ {len(issues)} hard issues — review needed")
        sys.exit(1)
    else:
        print(f"✓ No hard issues. {len(warnings)} soft warnings.")


if __name__ == "__main__":
    main()
